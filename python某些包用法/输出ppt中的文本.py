#!/usr/bin/python3
# coding: utf-8

# 输出ppt中的文本
from pptx import Presentation

path_to_presentation = '/home/gswewf/Downloads/诚信问题记录1.25.pptx'
prs = Presentation(path_to_presentation)

# text_runs will be populated with a list of strings,
# one for each text run in presentation
text_runs = []

for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                text_runs.append(run.text)

print(text_runs)

def main():
    pass


if __name__ == '__main__':
    main()